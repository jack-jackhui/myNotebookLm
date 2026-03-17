from typing import List
import os
import re
from urllib.parse import urlparse

from pptx import Presentation


def is_youtube_url(url: str) -> bool:
    """Check if the URL is a YouTube URL."""
    youtube_patterns = [
        r'(youtube\.com/watch\?v=)',
        r'(youtu\.be/)',
        r'(youtube\.com/embed/)',
        r'(youtube\.com/v/)',
    ]
    return any(re.search(pattern, url) for pattern in youtube_patterns)


def extract_youtube_video_id(url: str) -> str:
    """Extract video ID from YouTube URL."""
    patterns = [
        r'(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/embed/|youtube\.com/v/)([a-zA-Z0-9_-]{11})',
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    raise ValueError(f"Could not extract video ID from URL: {url}")


def extract_youtube_transcript(url: str) -> str:
    """Extract transcript from a YouTube video using youtube_transcript_api."""
    from youtube_transcript_api import YouTubeTranscriptApi

    video_id = extract_youtube_video_id(url)
    try:
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
        transcript_text = " ".join([entry['text'] for entry in transcript_list])
        return transcript_text
    except Exception as e:
        raise RuntimeError(f"Failed to extract YouTube transcript: {e}")


def extract_website_content(url: str) -> str:
    """Extract content from a website using newspaper3k."""
    from newspaper import Article

    try:
        article = Article(url)
        article.download()
        article.parse()

        content_parts = []
        if article.title:
            content_parts.append(article.title)
        if article.text:
            content_parts.append(article.text)

        return "\n\n".join(content_parts)
    except Exception as e:
        # Fallback to beautifulsoup4 for simpler extraction
        return extract_website_content_bs4(url)


def extract_website_content_bs4(url: str) -> str:
    """Fallback extraction using BeautifulSoup."""
    import requests
    from bs4 import BeautifulSoup

    response = requests.get(url, timeout=30)
    response.raise_for_status()

    soup = BeautifulSoup(response.text, 'html.parser')

    # Remove script and style elements
    for script in soup(["script", "style", "nav", "footer", "header"]):
        script.decompose()

    # Get text
    text = soup.get_text(separator='\n', strip=True)

    # Clean up whitespace
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    text = '\n'.join(chunk for chunk in chunks if chunk)

    return text


def extract_pdf_content(file_path: str) -> str:
    """Extract text content from a PDF file using PyMuPDF."""
    import fitz  # PyMuPDF

    text_parts = []
    with fitz.open(file_path) as doc:
        for page in doc:
            text_parts.append(page.get_text())

    return "\n".join(text_parts)


def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract text from PowerPoint file."""
    text = []
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def is_url(source: str) -> bool:
    """Check if source is a URL."""
    try:
        result = urlparse(source)
        return all([result.scheme, result.netloc])
    except:
        return False


def extract_content(source: str) -> str:
    """Extract content from a single source (URL or file path)."""
    ext = os.path.splitext(source)[1].lower()

    if ext == ".pptx":
        print(f"Extracting content from PowerPoint: {source}")
        return extract_text_from_pptx(source)
    elif ext == ".pdf":
        print(f"Extracting content from PDF: {source}")
        return extract_pdf_content(source)
    elif is_url(source):
        if is_youtube_url(source):
            print(f"Extracting transcript from YouTube: {source}")
            return extract_youtube_transcript(source)
        else:
            print(f"Extracting content from website: {source}")
            return extract_website_content(source)
    elif os.path.isfile(source):
        # Handle plain text files
        print(f"Reading text file: {source}")
        with open(source, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        raise ValueError(f"Unsupported source type: {source}")


def extract_content_from_sources(sources: List[str]) -> str:
    """Extract content from multiple sources and combine them."""
    extracted_texts = []
    for source in sources:
        extracted_text = extract_content(source)
        extracted_texts.append(extracted_text)
    combined_text = "\n\n".join(extracted_texts)
    return combined_text
